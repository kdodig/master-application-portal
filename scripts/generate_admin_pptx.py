#!/usr/bin/env python3
"""
Generates a non-technical feature presentation for the Admin area of the
application portal as a .pptx file with screenshot placeholders and speaker notes.

Requires: python-pptx
    pip install python-pptx

Output: slides/admin-portal-admin.pptx
"""

from datetime import datetime
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
except Exception as e:
    raise SystemExit(
        "python-pptx ist nicht installiert. Installiere es mit:\n"
        "    pip install python-pptx\n\n"
        f"Import-Fehler: {e}"
    )


OUTPUT_PATH = Path("slides/admin-portal-admin.pptx")


def add_notes(slide, text: str):
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    for i, line in enumerate(text.split("\n")):
        if i == 0:
            notes.text = line
        else:
            p = notes.add_paragraph()
            p.text = line


def add_bullets_textbox(slide, bullets, left=Inches(0.8), top=Inches(1.6), width=Inches(8.6), height=Inches(2.0)):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            tf.text = b
        else:
            p = tf.add_paragraph()
            p.text = b
        # style
        p = tf.paragraphs[-1]
        p.level = 0
        p.font.size = Pt(20)
    return textbox


def add_screenshot_placeholder(slide, label: str, left=Inches(0.8), top=Inches(3.0), width=Inches(8.6), height=Inches(3.9)):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    # No fill, dashed border, subtle label centered
    shape.fill.background()
    line = shape.line
    line.color.rgb = RGBColor(120, 120, 120)
    line.width = Pt(2)
    line.dash_style = MSO_LINE_DASH_STYLE.DASH

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"Screenshot hier einfügen: {label}"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(90, 90, 90)
    return shape


def build_deck():
    prs = Presentation()

    # Choose a simple theme default; layout indices:
    # 0=TITLE, 1=TITLE_AND_CONTENT, 5=TITLE_ONLY, 6=BLANK (varies by theme)
    TITLE_ONLY = 5

    slides = [
        {
            "title": "Admin-Bereich fürs Masterbewerbungsportal",
            "bullets": [
                "Kurzvorstellung für das bisherige Prüfteam",
                "Ziele: Transparenz, Effizienz, Nachvollziehbarkeit",
            ],
            "placeholder": "Titel / Key Visual",
            "notes": (
                "Zielgruppe sind Kolleg:innen, die bisher manuell geprüft haben.\n"
                "Fokus: Was kann der Admin-Bereich für euren Alltag leisten?"
            ),
        },
        {
            "title": "Ausgangssituation (bisher)",
            "bullets": [
                "Verteilte Dokumente, Medienbrüche, aufwendige Koordination",
                "Uneinheitliche Statusführung, wenig Auswertungen",
                "Zeitintensiv und fehleranfällig",
            ],
            "placeholder": "Beispielablauf (alt)",
            "notes": "Rahmen setzen, warum die Umstellung sinnvoll ist.",
        },
        {
            "title": "Ziele & Mehrwert",
            "bullets": [
                "Einheitlicher, geführter Workflow",
                "Schneller Überblick durch Filter und Kennzahlen",
                "Nachvollziehbarkeit: zentrale Daten, Timestamps, Status",
            ],
            "placeholder": "KPIs/Workflow-Diagramm",
            "notes": "Betont Nutzen für Review-Alltag und Auswertungen.",
        },
        {
            "title": "Login & Zugriff",
            "bullets": [
                "Geschützter Admin-Bereich mit Login",
                "Temporäre Passwörter erzwingen Änderung",
            ],
            "placeholder": "Login / Passwort ändern",
            "notes": "Kurz erklären: Zugriff nur für berechtigte Personen.",
        },
        {
            "title": "Dashboard & Statistiken",
            "bullets": [
                "Sofort-Überblick: Bewerbungen, Status, Länder",
                "Charts: Verteilung und Zeitverlauf",
            ],
            "placeholder": "Dashboard-Ansicht",
            "notes": "Hier sichtbar: zentrale KPIs für das Review-Team.",
        },
        {
            "title": "Bewerberübersicht",
            "bullets": [
                "Suche, Filter (z. B. Major, Status)",
                "Sortierung und Pagination",
                "Kontextmenü: direkte Aktionen/Review-Schritte",
            ],
            "placeholder": "Liste mit Filtern & Aktionen",
            "notes": "Zeigt, wie man schnell relevante Fälle findet.",
        },
        {
            "title": "Review: Dokumente",
            "bullets": [
                "Pflichtdokumente strukturiert prüfen und markieren",
                "Ablehnen oder fortfahren – klarer Status",
            ],
            "placeholder": "Dokumenten-Modal",
            "notes": "Hervorheben: Nachvollziehbare, einheitliche Dokumentenprüfung.",
        },
        {
            "title": "Review: Kursanalyse",
            "bullets": [
                "Kursdaten mit Vorschlägen prüfen und korrigieren",
                "PDF-Vorschau für Beschreibungen/Transcript",
            ],
            "placeholder": "Kursanalyse-Modal mit PDF-Viewer",
            "notes": "Idee: Unterstützt spätere Automatisierung, bleibt kontrollierbar.",
        },
        {
            "title": "Review: Personal Skills",
            "bullets": [
                "Skills mit Punkten erfassen und pflegen",
                "Validierung, klare Grenzen (Punktelimits)",
            ],
            "placeholder": "Skills-Modal",
            "notes": "Zeigt flexible, aber konsistente Erfassung.",
        },
        {
            "title": "Einstellungen: Accounts & Rollen",
            "bullets": [
                "Accounts verwalten: aktivieren/deaktivieren, bearbeiten",
                "Rollen zeigen Berechtigungen transparent",
            ],
            "placeholder": "Accounts-Übersicht mit Rollen",
            "notes": "Wer darf was? Einfache, nachvollziehbare Verwaltung.",
        },
        {
            "title": "Passwort ändern (Benutzer)",
            "bullets": [
                "Schnell und klar validiert",
                "Sicherheit durch verpflichtende Änderung bei Erstlogin",
            ],
            "placeholder": "Passwort-Ändern-Seite",
            "notes": "Hinweis auf Policies und Redirect zum Login.",
        },
        {
            "title": "Export & Berichte",
            "bullets": [
                "CSV-Export für schnelle Auswertungen",
                "Nutzung in Reporting/Controlling",
            ],
            "placeholder": "CSV-Export-Button",
            "notes": "Einfacher Transfer in gängige Tools (Excel, etc.).",
        },
        {
            "title": "Sicherheit & Datenschutz (Kurz)",
            "bullets": [
                "Authentifizierung & rollengestützter Zugriff",
                "Geordnete Dateispeicherung",
            ],
            "placeholder": "(Optionales Icon/Visual)",
            "notes": "Nur kurz – Detailtiefe bewusst gering halten.",
        },
        {
            "title": "Ausblick",
            "bullets": [
                "Benachrichtigungen/Erinnerungen",
                "Audit-Logs und feinere Rollen",
                "Weitere Automatisierungen",
            ],
            "placeholder": "Roadmap-Visual",
            "notes": "Was als Nächstes geplant ist – kein Commit.",
        },
        {
            "title": "Live-Demo Ablauf (optional)",
            "bullets": [
                "Login → Dashboard",
                "Bewerberliste: Filter & Aktionen",
                "Dokumente → Kurse → Skills",
                "Einstellungen & Export",
            ],
            "placeholder": "Ablauf-Checkliste",
            "notes": "Zeit im Blick behalten, Ladezustände abwarten.",
        },
    ]

    for slide_def in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[TITLE_ONLY])
        slide.shapes.title.text = slide_def["title"]
        add_bullets_textbox(slide, slide_def["bullets"])  # top area
        add_screenshot_placeholder(slide, slide_def["placeholder"])  # bottom area
        add_notes(slide, slide_def["notes"])  # speaker notes

    # Title slide first: move it to index 0 by creating a dedicated title slide
    title_slide = prs.slides._sldIdLst[0]
    # Add metadata: not strictly required, but nice to have
    prs.core_properties.author = "Admin-Team"
    prs.core_properties.title = "Admin-Bereich – Feature-Überblick"
    prs.core_properties.created = datetime.utcnow()

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    return OUTPUT_PATH


def main():
    out = build_deck()
    print(f"Präsentation erstellt: {out}")


if __name__ == "__main__":
    main()

